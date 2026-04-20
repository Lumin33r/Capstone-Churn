#!/usr/bin/env python3
"""Generate 'The Retention Engine' capstone presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette (dark tech theme) ────────────────────────────────────
BG_DARK    = RGBColor(0x0F, 0x17, 0x2A)   # deep navy
BG_CARD    = RGBColor(0x1A, 0x25, 0x3C)   # card bg
ACCENT     = RGBColor(0x00, 0xD4, 0xFF)   # cyan accent
ACCENT2    = RGBColor(0x7C, 0x3A, 0xED)   # purple accent
GREEN      = RGBColor(0x10, 0xB9, 0x81)   # green
ORANGE     = RGBColor(0xF5, 0x9E, 0x0B)   # amber
RED        = RGBColor(0xEF, 0x44, 0x44)   # red
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GRAY       = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)


def set_slide_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.shadow.inherit = False
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1.5)
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=LIGHT_GRAY, bullet_color=ACCENT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Segoe UI"
        p.space_after = Pt(6)
    return txBox


def section_divider(title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), fill_color=ACCENT)
    add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5), title, 48, ACCENT, True, PP_ALIGN.CENTER)
    if subtitle:
        add_text(slide, Inches(1), Inches(4.2), Inches(11), Inches(1), subtitle, 24, GRAY, False, PP_ALIGN.CENTER)
    return slide


def dev_slide(name, role, emoji, contributions, tech_stack, talking_points):
    """Create an individual developer slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    # top accent bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=ACCENT2)
    # developer name + role header
    add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
             f"{emoji}  {name}  —  {role}", 36, WHITE, True, PP_ALIGN.LEFT)
    # separator line
    add_shape(slide, Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.03), fill_color=ACCENT)

    # LEFT COLUMN — What I Built
    add_shape(slide, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.4), fill_color=BG_CARD, line_color=RGBColor(0x33, 0x44, 0x66))
    add_text(slide, Inches(0.9), Inches(1.8), Inches(5), Inches(0.5), "WHAT I BUILT", 20, ACCENT, True)
    add_bullet_list(slide, Inches(0.9), Inches(2.4), Inches(5.2), Inches(4.2), contributions, 15)

    # RIGHT COLUMN — Tech Stack + Talking Points
    add_shape(slide, Inches(6.8), Inches(1.6), Inches(5.8), Inches(2.4), fill_color=BG_CARD, line_color=RGBColor(0x33, 0x44, 0x66))
    add_text(slide, Inches(7.1), Inches(1.8), Inches(5), Inches(0.5), "TECH STACK", 20, GREEN, True)
    add_bullet_list(slide, Inches(7.1), Inches(2.4), Inches(5.2), Inches(1.8), tech_stack, 14, GRAY)

    add_shape(slide, Inches(6.8), Inches(4.2), Inches(5.8), Inches(2.8), fill_color=BG_CARD, line_color=RGBColor(0x33, 0x44, 0x66))
    add_text(slide, Inches(7.1), Inches(4.4), Inches(5), Inches(0.5), "KEY TALKING POINTS", 20, ORANGE, True)
    add_bullet_list(slide, Inches(7.1), Inches(5.0), Inches(5.2), Inches(2.2), talking_points, 14, GRAY)

    return slide


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), fill_color=ACCENT)
add_text(slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
         "THE RETENTION ENGINE", 60, WHITE, True, PP_ALIGN.CENTER, "Segoe UI")
add_text(slide, Inches(1), Inches(3.2), Inches(11.3), Inches(1),
         "AI-Powered Customer Churn Prediction & Retention Platform", 28, ACCENT, False, PP_ALIGN.CENTER)
add_shape(slide, Inches(4), Inches(4.5), Inches(5.3), Inches(0.03), fill_color=GRAY)
add_text(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.6),
         "Troy  ·  Kathleen  ·  Okino  ·  George", 24, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.7), Inches(11.3), Inches(0.5),
         "Code Platoon  ·  DevOps & Cloud Engineering  ·  April 2026", 18, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=RED)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         "THE PROBLEM", 40, RED, True, PP_ALIGN.LEFT)

# stat cards
stats = [
    ("$1.6T", "lost annually to\ncustomer churn", RED),
    ("5–25×", "more expensive to\nacquire vs. retain", ORANGE),
    ("67%", "of churn is\npreventable", GREEN),
]
for i, (big, desc, clr) in enumerate(stats):
    x = Inches(0.8 + i * 4.2)
    add_shape(slide, x, Inches(1.8), Inches(3.6), Inches(2.2), fill_color=BG_CARD, line_color=clr)
    add_text(slide, x + Inches(0.3), Inches(2.0), Inches(3), Inches(1), big, 52, clr, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(3.0), Inches(3), Inches(0.8), desc, 18, GRAY, False, PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(1.5),
         "Call center managers fly blind — no real-time visibility into which customers\n"
         "are about to leave, no sentiment awareness from support calls, and no AI-driven\n"
         "retention playbook. Decisions are reactive, not predictive.",
         20, LIGHT_GRAY, False, PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=GREEN)
add_text(slide, Inches(0.8), Inches(0.4), Inches(12), Inches(0.8),
         "THE SOLUTION — Manager's Command Center", 36, GREEN, True)

features = [
    ("🎯  ANALYZE", "Select a customer, optionally paste a call transcript.\nGet churn risk score, sentiment analysis, and AI-generated retention actions.", ACCENT),
    ("💬  CHAT", "Conversational AI powered by Claude Haiku 4.5 via Bedrock.\nAsk about any customer, get high-risk reports, or strategize retention.", ACCENT2),
    ("📞  TRANSCRIBE", "Upload call audio → Amazon Transcribe with speaker diarization.\nAutomatic sentiment enrichment feeds directly into churn predictions.", GREEN),
]
for i, (title, desc, clr) in enumerate(features):
    y = Inches(1.5 + i * 1.9)
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(1.7), fill_color=BG_CARD, line_color=clr)
    add_text(slide, Inches(1.0), y + Inches(0.15), Inches(4), Inches(0.5), title, 24, clr, True)
    add_text(slide, Inches(1.0), y + Inches(0.65), Inches(11), Inches(1), desc, 16, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=ACCENT)
add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
         "PLATFORM ARCHITECTURE", 36, ACCENT, True)

# Architecture layers (simplified visual)
layers = [
    ("FRONTEND", "React + Vite + Tailwind  ·  3 Tabs: Analyze | Chat | Transcribe", Inches(1.2), ACCENT),
    ("AGENT SERVICE", "LangGraph Orchestrator  ·  Claude Haiku 4.5  ·  Bedrock Guardrails  ·  5 Tools", Inches(2.5), ACCENT2),
    ("ML SERVICES", "Churn Predictor API (:8001)  ·  Sentiment Analysis API (:8002)", Inches(3.8), GREEN),
    ("AWS SERVICES", "SageMaker Endpoints (XGBoost + Sentiment)  ·  Bedrock  ·  S3  ·  Transcribe", Inches(5.1), ORANGE),
    ("INFRASTRUCTURE", "Terraform  ·  EKS (Kubernetes)  ·  GitHub Actions CI/CD  ·  GHCR  ·  Slack", Inches(6.3), GRAY),
]
for label, desc, y, clr in layers:
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(1.1), fill_color=BG_CARD, line_color=clr)
    add_text(slide, Inches(1.0), y + Inches(0.1), Inches(3), Inches(0.5), label, 20, clr, True)
    add_text(slide, Inches(4.2), y + Inches(0.15), Inches(8), Inches(0.8), desc, 16, LIGHT_GRAY)

# arrows between layers
for y_top in [Inches(2.3), Inches(3.6), Inches(4.9), Inches(6.1)]:
    add_text(slide, Inches(6.3), y_top, Inches(0.8), Inches(0.3), "▼", 18, GRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=ACCENT)
add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
         "TECH STACK", 36, ACCENT, True)

cols = [
    ("ML / AI", [
        "Amazon SageMaker (XGBoost)",
        "Amazon Bedrock (Claude Haiku 4.5)",
        "Bedrock Guardrails (PII, Content)",
        "Amazon Transcribe (Speaker Diarization)",
        "LangGraph + LangChain Agents",
        "LangSmith Tracing",
    ], ACCENT2),
    ("BACKEND", [
        "FastAPI (3 microservices)",
        "Python 3.12",
        "boto3 / httpx / pydantic",
        "Docker (multi-stage builds)",
        "S3 data storage (parquet/CSV)",
    ], GREEN),
    ("FRONTEND", [
        "React 18 + TypeScript",
        "Vite + Tailwind CSS",
        "Lucide React Icons",
        "nginx (production)",
    ], ACCENT),
    ("INFRA / DEVOPS", [
        "Terraform (IaC)",
        "Amazon EKS (Kubernetes)",
        "GitHub Actions (6 workflows)",
        "GHCR (container registry)",
        "Slack notifications",
    ], ORANGE),
]
for i, (title, items, clr) in enumerate(cols):
    x = Inches(0.4 + i * 3.25)
    add_shape(slide, x, Inches(1.3), Inches(3.0), Inches(5.8), fill_color=BG_CARD, line_color=clr)
    add_text(slide, x + Inches(0.2), Inches(1.5), Inches(2.6), Inches(0.5), title, 20, clr, True, PP_ALIGN.CENTER)
    add_bullet_list(slide, x + Inches(0.2), Inches(2.2), Inches(2.6), Inches(4.5), items, 13, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — LIVE DEMO  (placeholder)
# ═══════════════════════════════════════════════════════════════════════
slide = section_divider("LIVE DEMO", "Manager's Command Center  ·  Analyze  ·  Chat  ·  Transcribe")


# ═══════════════════════════════════════════════════════════════════════
# DEVELOPER SECTION DIVIDER
# ═══════════════════════════════════════════════════════════════════════
section_divider("MEET THE TEAM", "Individual Contributions")


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — TROY
# ═══════════════════════════════════════════════════════════════════════
dev_slide(
    name="Troy",
    role="CI/CD & Project Management",
    emoji="🔧",
    contributions=[
        "Designed and implemented all 6 GitHub Actions workflows",
        "terraform.yml — plan/apply with environment selection (dev/staging/prod)",
        "sagemaker-deploy.yml — endpoint lifecycle with health polling + inference tests",
        "deploy.yml — matrix Docker build → GHCR → EKS rollout + smoke tests",
        "ci-post-merge.yml — automated 4-job pipeline on merge to main",
        "slack-pr-events.yml — team notifications on PR open/merge",
        "teardown.yml — full infrastructure teardown workflow",
        "Set up GitHub Projects Kanban board for sprint tracking",
        "Branch protection rules and PR review workflow",
    ],
    tech_stack=[
        "GitHub Actions (YAML, matrix builds, workflow_dispatch)",
        "Docker & GHCR (container registry)",
        "kubectl (EKS rollout management)",
        "Slack Webhooks",
        "GitHub Projects (Kanban)",
    ],
    talking_points=[
        "Walk through the CI/CD pipeline diagram",
        "Show the manual dispatch → automatic validation flow",
        "Demonstrate Slack notification on PR merge",
        "Explain how the post-merge pipeline catches regressions",
    ]
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — KATHLEEN
# ═══════════════════════════════════════════════════════════════════════
dev_slide(
    name="Kathleen",
    role="ML Engineering — Churn Prediction & Frontend",
    emoji="📊",
    contributions=[
        "Trained XGBoost churn model v3: 95% accuracy, 0.9861 AUC",
        "Engineered 31 features including 7 cross-model Agent 1 features",
        "Deployed model to SageMaker endpoint (XGBoost container, native format)",
        "Built churn-predictor-api with internal S3 customer data lookup",
        "/predict, /customers, /customer-details, /high-risk endpoints",
        "Co-built LangGraph retention_agent.py with 4-tool orchestration",
        "Designed TriLink product catalog with risk-based retention actions",
        "Implemented output guardrails for agent recommendations",
        "Built React + Tailwind frontend (Analyze, Chat, Transcribe tabs)",
        "Multi-stage Dockerfile (node build → nginx serve)",
    ],
    tech_stack=[
        "SageMaker (XGBoost, endpoint deployment)",
        "FastAPI + boto3 + pandas",
        "LangGraph + LangChain (agent orchestration)",
        "React 18 + TypeScript + Vite + Tailwind",
        "Docker (multi-stage builds)",
    ],
    talking_points=[
        "Walk through the churn model training + feature engineering",
        "Show the 31-feature pipeline including Agent 1 enrichment",
        "Demo the /high-risk batch prediction endpoint",
        "Explain the LangGraph agent routing logic + guardrails",
    ]
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — OKINO
# ═══════════════════════════════════════════════════════════════════════
dev_slide(
    name="Okino",
    role="ML Engineering — Sentiment Analysis",
    emoji="🧠",
    contributions=[
        "Built sentiment-analysis-api with Bedrock-powered analysis",
        "Training notebook for QA/transcript evaluation model",
        "FastAPI wrapper returning qa_score, sentiment, emotion metrics",
        "Designed the 7 Agent 1 features that enrich churn predictions:",
        "  → sentiment, emotion_frustration, emotion_anger",
        "  → sentiment_shift, escalation_flag, resolution_flag, qa_score",
        "Co-built LangGraph agent tool contracts for cross-service calls",
        "Dockerfile + K8s deployment and service manifests",
        "SageMaker endpoint configuration in Terraform",
    ],
    tech_stack=[
        "SageMaker (model training + endpoints)",
        "Amazon Bedrock (Claude for sentiment analysis)",
        "FastAPI + boto3",
        "Docker + Kubernetes manifests",
        "Terraform (sagemaker.tf)",
    ],
    talking_points=[
        "Explain how call transcripts become 7 structured features",
        "Show the sentiment → churn enrichment pipeline",
        "Walk through the /predict response schema",
        "Describe how Agent 1 features improve churn accuracy",
    ]
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — GEORGE
# ═══════════════════════════════════════════════════════════════════════
dev_slide(
    name="George",
    role="Infrastructure & Platform Engineering",
    emoji="⚙️",
    contributions=[
        "Authored all Terraform IaC: S3, IAM, SageMaker, EKS, Bedrock, Transcribe",
        "Designed & deployed full Kubernetes architecture (EKS cluster)",
        "All K8s manifests: deployments, services, configmaps, secrets, quotas",
        "Built agent-service: FastAPI + LangGraph /chat route with session memory",
        "Integrated Bedrock Guardrails (PII anonymization, content filtering)",
        "docker-compose.yml for full local development stack",
        "Amazon Transcribe pipeline: S3 audio → Lambda → Transcribe → S3 transcripts",
        "Terraform for Lambda + IAM + S3 event trigger",
    ],
    tech_stack=[
        "Terraform (14 config files, multi-resource)",
        "Amazon EKS + kubectl",
        "Kubernetes (namespace, RBAC, resource quotas)",
        "FastAPI + LangGraph + Bedrock",
        "AWS Lambda + Amazon Transcribe",
        "Docker + docker-compose",
    ],
    talking_points=[
        "Walk through the Terraform resource graph",
        "Show the K8s namespace layout and service mesh",
        "Explain the dual-LLM guardrail architecture",
        "Demo the Transcribe pipeline: audio → structured transcript",
    ]
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — CI/CD PIPELINE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=ACCENT)
add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
         "CI/CD PIPELINE — 6 GITHUB ACTIONS WORKFLOWS", 32, ACCENT, True)

workflows = [
    ("terraform.yml", "Manual Dispatch", "Terraform plan → apply  (dev/staging/prod)", ACCENT2),
    ("sagemaker-deploy.yml", "Manual Dispatch", "Deploy/delete endpoints + health poll + inference test", GREEN),
    ("deploy.yml", "Manual Dispatch", "Matrix Docker build → GHCR → EKS rollout + smoke test", ACCENT),
    ("ci-post-merge.yml", "Push to main", "Unit tests → Docker → Endpoint health → E2E smoke", ORANGE),
    ("slack-pr-events.yml", "PR open/merge", "Formatted Slack notifications to team channel", RGBColor(0xE8, 0x79, 0xF9)),
    ("teardown.yml", "Manual Dispatch", "Full infrastructure teardown", RED),
]
for i, (name, trigger, desc, clr) in enumerate(workflows):
    y = Inches(1.2 + i * 1.0)
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(0.85), fill_color=BG_CARD, line_color=clr)
    add_text(slide, Inches(1.0), y + Inches(0.1), Inches(3.5), Inches(0.35), name, 17, clr, True)
    add_text(slide, Inches(1.0), y + Inches(0.45), Inches(3.5), Inches(0.35), trigger, 13, GRAY)
    add_text(slide, Inches(4.8), y + Inches(0.2), Inches(7.5), Inches(0.5), desc, 16, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — AGENTIC AI DEEP DIVE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=ACCENT2)
add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
         "AGENTIC AI — LangGraph Orchestration", 36, ACCENT2, True)

# Agent nodes
nodes = [
    ("DataGatherer", "Bedrock + Guardrails\nCollects customer data, runs tools", ACCENT, Inches(1.5)),
    ("RetentionStrategist", "Clean LLM (no guardrail)\nGenerates retention offers\nfrom risk analysis", ACCENT2, Inches(3.3)),
]
for name, desc, clr, y in nodes:
    add_shape(slide, Inches(0.6), y, Inches(5.5), Inches(1.5), fill_color=BG_CARD, line_color=clr)
    add_text(slide, Inches(1.0), y + Inches(0.15), Inches(4.8), Inches(0.4), name, 22, clr, True)
    add_text(slide, Inches(1.0), y + Inches(0.6), Inches(4.8), Inches(0.8), desc, 14, GRAY)

# Tools list
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(3.3), fill_color=BG_CARD, line_color=GREEN)
add_text(slide, Inches(7.1), Inches(1.7), Inches(5.2), Inches(0.4), "5 AGENT TOOLS", 22, GREEN, True)
tools = [
    "get_customer_details → Churn API /customer-details/{id}",
    "predict_churn → Churn API /predict (31 features → SageMaker)",
    "analyze_call → Sentiment API /predict (transcript → 7 metrics)",
    "get_high_risk_customers → Churn API /high-risk (batch)",
    "get_transcripts → S3 transcript listing per customer",
]
add_bullet_list(slide, Inches(7.1), Inches(2.3), Inches(5.2), Inches(2.5), tools, 13, LIGHT_GRAY)

# Guardrails callout
add_shape(slide, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.8), fill_color=BG_CARD, line_color=ORANGE)
add_text(slide, Inches(1.0), Inches(5.4), Inches(5), Inches(0.4), "BEDROCK GUARDRAILS", 20, ORANGE, True)
guardrails = [
    "Hate/Sexual content: HIGH filter  ·  Violence: HIGH/MEDIUM  ·  Misconduct: HIGH/LOW",
    "PII Protection: ANONYMIZE names, emails, phone, SSN  ·  BLOCK credit card, bank account",
    "Dual-LLM design: DataGatherer has guardrails, RetentionStrategist uses clean LLM to avoid blocking domain terms (cancel, churn)"
]
add_bullet_list(slide, Inches(1.0), Inches(5.9), Inches(11.5), Inches(1.5), guardrails, 13, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — KEY METRICS & RESULTS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=GREEN)
add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
         "KEY METRICS & RESULTS", 36, GREEN, True)

metrics = [
    ("95%", "Model\nAccuracy", GREEN),
    ("0.986", "AUC\nScore", ACCENT),
    ("31", "Engineered\nFeatures", ACCENT2),
    ("5", "Agent\nTools", ORANGE),
    ("3", "Microservices\n+ Frontend", RGBColor(0xE8, 0x79, 0xF9)),
    ("6", "CI/CD\nWorkflows", RED),
]
for i, (val, label, clr) in enumerate(metrics):
    x = Inches(0.5 + i * 2.15)
    add_shape(slide, x, Inches(1.5), Inches(1.9), Inches(2.2), fill_color=BG_CARD, line_color=clr)
    add_text(slide, x + Inches(0.1), Inches(1.7), Inches(1.7), Inches(1), val, 44, clr, True, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), Inches(2.7), Inches(1.7), Inches(0.8), label, 14, GRAY, False, PP_ALIGN.CENTER)

# highlights
highlights = [
    "End-to-end ML pipeline: raw data → training → SageMaker endpoint → FastAPI wrapper → agent tool → user response",
    "Agentic AI with multi-step reasoning: transcript → sentiment → churn enrichment → retention strategy",
    "Full IaC: 14 Terraform files provisioning S3, IAM, SageMaker, EKS, Bedrock, Transcribe, Lambda",
    "Production K8s: namespace isolation, resource quotas, health probes, ConfigMaps, rolling deploys",
    "Amazon Transcribe integration: audio upload → speaker diarization → structured transcript → S3",
]
add_shape(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(3.2), fill_color=BG_CARD, line_color=ACCENT)
add_text(slide, Inches(0.8), Inches(4.2), Inches(5), Inches(0.4), "PLATFORM HIGHLIGHTS", 20, ACCENT, True)
add_bullet_list(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(2.8), highlights, 15, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — LESSONS LEARNED
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=ACCENT2)
add_text(slide, Inches(0.8), Inches(0.3), Inches(12), Inches(0.7),
         "LESSONS LEARNED", 36, ACCENT2, True)

lessons_left = [
    "Define API contracts early — enabled parallel development",
    "Bedrock Guardrails block domain terms (cancel, churn) — needed dual-LLM workaround",
    "SageMaker endpoint cold starts affect initial latency — added health polling in CI",
    "Multi-stage Docker builds cut image sizes significantly",
]
lessons_right = [
    "LangGraph > raw LangChain for multi-step agent workflows",
    "docker-compose essential for local integration testing before K8s",
    "Slack notifications keep the entire team informed on PR flow",
    "Feature engineering (31 features) matters more than model choice",
]

add_shape(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(5.5), fill_color=BG_CARD, line_color=ACCENT)
add_text(slide, Inches(0.9), Inches(1.5), Inches(5.2), Inches(0.4), "TECHNICAL INSIGHTS", 20, ACCENT, True)
add_bullet_list(slide, Inches(0.9), Inches(2.1), Inches(5.2), Inches(4.5), lessons_left, 15, LIGHT_GRAY)

add_shape(slide, Inches(6.8), Inches(1.3), Inches(5.8), Inches(5.5), fill_color=BG_CARD, line_color=GREEN)
add_text(slide, Inches(7.1), Inches(1.5), Inches(5.2), Inches(0.4), "PROCESS INSIGHTS", 20, GREEN, True)
add_bullet_list(slide, Inches(7.1), Inches(2.1), Inches(5.2), Inches(4.5), lessons_right, 15, LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — Q&A / THANK YOU
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.1), fill_color=ACCENT)
add_text(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.5),
         "QUESTIONS?", 60, WHITE, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(4), Inches(3.8), Inches(5.3), Inches(0.03), fill_color=ACCENT)
add_text(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.8),
         "THE RETENTION ENGINE", 28, ACCENT, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.0), Inches(11.3), Inches(0.5),
         "Troy  ·  Kathleen  ·  Okino  ·  George", 22, LIGHT_GRAY, False, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(5.8), Inches(11.3), Inches(0.5),
         "github.com/Lumin33r/Capstone-Churn", 18, GRAY, False, PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(__file__), "TheRetentionEngine.pptx")
prs.save(out_path)
print(f"✅ Saved → {out_path}")
