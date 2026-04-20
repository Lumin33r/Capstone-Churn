#!/usr/bin/env python3
"""Generate 'The Retention Engine' presentation — Norman Rockwell theme.

Warm Americana palette, serif fonts, parchment backgrounds, vintage card styling.
Image placeholders describe Rockwell-style scenes to insert manually.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Norman Rockwell color palette ──────────────────────────────────────
PARCHMENT   = RGBColor(0xF5, 0xF0, 0xE1)   # warm cream background
PARCHMENT2  = RGBColor(0xEA, 0xE0, 0xCC)   # darker parchment for cards
CARD_BG     = RGBColor(0xFD, 0xF8, 0xEE)   # light card bg
DEEP_RED    = RGBColor(0x8B, 0x1A, 0x1A)   # Rockwell red
NAVY        = RGBColor(0x1B, 0x2A, 0x4A)   # deep navy blue
WARM_GOLD   = RGBColor(0xC4, 0x8A, 0x2A)   # warm gold
FOREST      = RGBColor(0x2D, 0x5A, 0x27)   # forest green
BROWN       = RGBColor(0x5C, 0x3D, 0x2E)   # warm brown
DARK_BROWN  = RGBColor(0x3E, 0x28, 0x1E)   # dark brown text
CHARCOAL    = RGBColor(0x33, 0x33, 0x33)   # dark text
SEPIA       = RGBColor(0x70, 0x4A, 0x2A)   # sepia accent
BURNT_ORANGE = RGBColor(0xBF, 0x5B, 0x16)  # warm orange
SLATE       = RGBColor(0x64, 0x5C, 0x54)   # warm gray
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
CREAM_TEXT  = RGBColor(0xF9, 0xF5, 0xEB)

FONT = "Georgia"
FONT_SANS = "Garamond"


def set_slide_bg(slide, color=PARCHMENT):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None,
              shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_width=Pt(1.5)):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.shadow.inherit = False
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    else:
        shp.fill.background()
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = line_width
    else:
        shp.line.fill.background()
    return shp


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(1.5)):
    return add_shape(slide, left, top, width, height, fill_color, line_color,
                     MSO_SHAPE.RECTANGLE, line_width)


def add_text(slide, left, top, width, height, text, font_size=18, color=CHARCOAL,
             bold=False, alignment=PP_ALIGN.LEFT, font_name=FONT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_BROWN, bullet="✦"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(6)
    return txBox


def add_image_placeholder(slide, left, top, width, height, description, border_color=BROWN):
    """Add a vintage-framed placeholder with scene description."""
    # outer frame
    add_rect(slide, left, top, width, height, fill_color=RGBColor(0xD4, 0xC8, 0xAE),
             line_color=border_color, line_width=Pt(3))
    # inner frame
    margin = Inches(0.12)
    add_rect(slide, left + margin, top + margin,
             width - margin * 2, height - margin * 2,
             fill_color=RGBColor(0xE8, 0xDD, 0xC8),
             line_color=RGBColor(0xA0, 0x8C, 0x72), line_width=Pt(1))
    # description text
    add_text(slide, left + Inches(0.25), top + Inches(0.25),
             width - Inches(0.5), height - Inches(0.5),
             f"🖼  {description}", 12, SEPIA, False, PP_ALIGN.CENTER, FONT, italic=True)


def ornamental_rule(slide, left, top, width):
    """A decorative horizontal rule."""
    add_text(slide, left, top, width, Inches(0.3),
             "— ✦ — ✦ — ✦ —", 16, WARM_GOLD, False, PP_ALIGN.CENTER, FONT)


def top_banner(slide, color=DEEP_RED, height=Inches(1.1)):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), height,
             fill_color=color, line_color=None)
    # thin gold line at bottom of banner
    add_rect(slide, Inches(0), height - Inches(0.04), Inches(13.333), Inches(0.04),
             fill_color=WARM_GOLD, line_color=None)


def section_divider(title, subtitle="", img_desc=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PARCHMENT)
    top_banner(slide, NAVY)
    add_text(slide, Inches(1), Inches(0.15), Inches(11.3), Inches(0.8),
             title, 44, CREAM_TEXT, True, PP_ALIGN.CENTER, FONT)
    ornamental_rule(slide, Inches(3), Inches(2.8), Inches(7.3))
    if subtitle:
        add_text(slide, Inches(1), Inches(3.3), Inches(11.3), Inches(0.8),
                 subtitle, 24, SEPIA, False, PP_ALIGN.CENTER, FONT, italic=True)
    if img_desc:
        add_image_placeholder(slide, Inches(4.5), Inches(4.3), Inches(4.3), Inches(2.8), img_desc, NAVY)
    return slide


def dev_slide(name, role, emoji, contributions, tech_stack, talking_points, img_desc, accent=DEEP_RED):
    """Individual developer slide, vintage style."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PARCHMENT)
    top_banner(slide, accent, Inches(0.9))
    add_text(slide, Inches(0.8), Inches(0.1), Inches(11.5), Inches(0.7),
             f"{name}  —  {role}", 34, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

    # Portrait placeholder (left)
    add_image_placeholder(slide, Inches(0.6), Inches(1.2), Inches(3.2), Inches(3.0),
                          img_desc, accent)

    # What I Built card (right of portrait)
    add_shape(slide, Inches(4.1), Inches(1.2), Inches(8.6), Inches(3.0),
              fill_color=CARD_BG, line_color=BROWN)
    add_text(slide, Inches(4.4), Inches(1.35), Inches(5), Inches(0.4),
             "What I Built", 22, accent, True, PP_ALIGN.LEFT, FONT)
    add_bullet_list(slide, Inches(4.4), Inches(1.85), Inches(8), Inches(2.2),
                    contributions, 13, DARK_BROWN, "✦")

    # Bottom row — Tech Stack (left) + Talking Points (right)
    add_shape(slide, Inches(0.6), Inches(4.5), Inches(5.8), Inches(2.7),
              fill_color=CARD_BG, line_color=FOREST)
    add_text(slide, Inches(0.9), Inches(4.65), Inches(5), Inches(0.4),
             "Tech Stack", 20, FOREST, True, PP_ALIGN.LEFT, FONT)
    add_bullet_list(slide, Inches(0.9), Inches(5.15), Inches(5.2), Inches(2.0),
                    tech_stack, 13, DARK_BROWN, "◈")

    add_shape(slide, Inches(6.8), Inches(4.5), Inches(5.9), Inches(2.7),
              fill_color=CARD_BG, line_color=WARM_GOLD)
    add_text(slide, Inches(7.1), Inches(4.65), Inches(5), Inches(0.4),
             "Key Talking Points", 20, WARM_GOLD, True, PP_ALIGN.LEFT, FONT)
    add_bullet_list(slide, Inches(7.1), Inches(5.15), Inches(5.3), Inches(2.0),
                    talking_points, 13, DARK_BROWN, "◈")

    return slide


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
# Large navy banner
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(2.0),
         fill_color=NAVY, line_color=None)
add_rect(slide, Inches(0), Inches(1.96), Inches(13.333), Inches(0.06),
         fill_color=WARM_GOLD, line_color=None)

add_text(slide, Inches(1), Inches(0.3), Inches(11.3), Inches(1.3),
         "THE RETENTION ENGINE", 62, CREAM_TEXT, True, PP_ALIGN.CENTER, FONT)

add_text(slide, Inches(1), Inches(1.3), Inches(11.3), Inches(0.6),
         "An AI-Powered Customer Churn Prediction & Retention Platform",
         24, WARM_GOLD, False, PP_ALIGN.CENTER, FONT, italic=True)

# Central image placeholder
add_image_placeholder(slide, Inches(3.8), Inches(2.5), Inches(5.7), Inches(3.2),
    "Norman Rockwell-style painting: A friendly telephone operator at a switchboard, "
    "connecting calls with a warm smile — representing customer connection & retention",
    NAVY)

ornamental_rule(slide, Inches(2.5), Inches(5.9), Inches(8.3))

add_text(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.5),
         "Troy  ·  Kathleen  ·  Okino  ·  George", 26, NAVY, True, PP_ALIGN.CENTER, FONT)
add_text(slide, Inches(1), Inches(6.7), Inches(11.3), Inches(0.4),
         "Code Platoon  ·  DevOps & Cloud Engineering  ·  April 2026", 16, SEPIA, False, PP_ALIGN.CENTER, FONT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, DEEP_RED)
add_text(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7),
         "THE PROBLEM", 40, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

# Image placeholder — left side
add_image_placeholder(slide, Inches(0.6), Inches(1.4), Inches(4.0), Inches(3.5),
    "Rockwell-style: A worried shopkeeper looking at an empty store, "
    "a 'Going Out of Business' sign visible — customer loss & churn",
    DEEP_RED)

# Stat cards — right side
stats = [
    ("$1.6 Trillion", "lost annually to customer churn", DEEP_RED),
    ("5–25× More Expensive", "to acquire a new customer\nthan to keep one", BURNT_ORANGE),
    ("67% Preventable", "with the right tools\nand early intervention", FOREST),
]
for i, (big, desc, clr) in enumerate(stats):
    y = Inches(1.4 + i * 1.2)
    add_shape(slide, Inches(5.0), y, Inches(7.7), Inches(1.05),
              fill_color=CARD_BG, line_color=clr)
    add_text(slide, Inches(5.3), y + Inches(0.08), Inches(4), Inches(0.45), big, 22, clr, True, font_name=FONT)
    add_text(slide, Inches(5.3), y + Inches(0.5), Inches(7), Inches(0.5), desc, 14, SLATE, font_name=FONT)

add_text(slide, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.5),
         "Call center managers fly blind — no real-time visibility into which customers "
         "are about to leave, no sentiment awareness from support calls, and no AI-driven "
         "retention playbook. Decisions are reactive, not predictive.",
         18, DARK_BROWN, False, PP_ALIGN.LEFT, FONT, italic=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, FOREST)
add_text(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7),
         "THE SOLUTION — Manager's Command Center", 34, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

features = [
    ("Analyze", "Select a customer, optionally paste a call transcript. "
     "Get churn risk score, sentiment analysis, and AI-generated retention actions.",
     DEEP_RED,
     "Rockwell-style: A doctor carefully examining a patient's chart with a magnifying glass — "
     "representing careful analysis of customer health"),
    ("Chat", "Conversational AI powered by Claude Haiku 4.5 via Amazon Bedrock. "
     "Ask about any customer, get high-risk reports, or strategize retention.",
     NAVY,
     "Rockwell-style: Two colleagues chatting warmly over a desk, one pointing at notes — "
     "representing the AI conversation partner"),
    ("Transcribe", "Upload call audio → Amazon Transcribe with speaker diarization. "
     "Automatic sentiment enrichment feeds directly into churn predictions.",
     FOREST,
     "Rockwell-style: An old-time stenographer carefully transcribing notes from a phonograph — "
     "representing audio-to-text transformation"),
]
for i, (title, desc, clr, img) in enumerate(features):
    y = Inches(1.3 + i * 2.05)
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(1.85),
              fill_color=CARD_BG, line_color=clr)
    add_image_placeholder(slide, Inches(0.8), y + Inches(0.15), Inches(2.2), Inches(1.55), img, clr)
    add_text(slide, Inches(3.2), y + Inches(0.15), Inches(4), Inches(0.45), title, 26, clr, True, font_name=FONT)
    add_text(slide, Inches(3.2), y + Inches(0.65), Inches(9.2), Inches(1.0), desc, 15, DARK_BROWN, font_name=FONT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE OVERVIEW
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, NAVY)
add_text(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7),
         "PLATFORM ARCHITECTURE", 36, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

layers = [
    ("Frontend", "React + Vite + Tailwind  ·  3 Tabs: Analyze | Chat | Transcribe", DEEP_RED),
    ("Agent Service", "LangGraph Orchestrator  ·  Claude Haiku 4.5  ·  Bedrock Guardrails  ·  5 Tools", NAVY),
    ("ML Services", "Churn Predictor API (:8001)  ·  Sentiment Analysis API (:8002)", FOREST),
    ("AWS Services", "SageMaker Endpoints (XGBoost + Sentiment)  ·  Bedrock  ·  S3  ·  Transcribe", WARM_GOLD),
    ("Infrastructure", "Terraform  ·  EKS (Kubernetes)  ·  GitHub Actions CI/CD  ·  GHCR  ·  Slack", BROWN),
]
for i, (label, desc, clr) in enumerate(layers):
    y = Inches(1.25 + i * 1.15)
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(1.0),
              fill_color=CARD_BG, line_color=clr)
    add_text(slide, Inches(1.0), y + Inches(0.15), Inches(3), Inches(0.4), label, 22, clr, True, font_name=FONT)
    add_text(slide, Inches(4.2), y + Inches(0.2), Inches(8), Inches(0.6), desc, 16, DARK_BROWN, font_name=FONT)

# arrows
for y_top in [Inches(2.25), Inches(3.4), Inches(4.55), Inches(5.7)]:
    add_text(slide, Inches(6.3), y_top, Inches(0.8), Inches(0.3), "▼", 18, SEPIA, False, PP_ALIGN.CENTER)

# Side illustration
add_image_placeholder(slide, Inches(9.5), Inches(1.25), Inches(3.2), Inches(2.4),
    "Rockwell-style: A craftsman carefully assembling layers of a wooden cabinet — "
    "representing the layered architecture",
    NAVY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECH STACK
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, BROWN)
add_text(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7),
         "TECH STACK — Our Toolkit", 36, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

cols = [
    ("ML / AI", [
        "Amazon SageMaker (XGBoost)",
        "Amazon Bedrock (Claude Haiku 4.5)",
        "Bedrock Guardrails (PII, Content)",
        "Amazon Transcribe",
        "LangGraph + LangChain Agents",
        "LangSmith Tracing",
    ], DEEP_RED),
    ("Backend", [
        "FastAPI (3 microservices)",
        "Python 3.12",
        "boto3 / httpx / pydantic",
        "Docker (multi-stage builds)",
        "S3 data storage (parquet/CSV)",
    ], NAVY),
    ("Frontend", [
        "React 18 + TypeScript",
        "Vite + Tailwind CSS",
        "Lucide React Icons",
        "nginx (production)",
    ], FOREST),
    ("Infra / DevOps", [
        "Terraform (IaC)",
        "Amazon EKS (Kubernetes)",
        "GitHub Actions (6 workflows)",
        "GHCR (container registry)",
        "Slack notifications",
    ], WARM_GOLD),
]
for i, (title, items, clr) in enumerate(cols):
    x = Inches(0.4 + i * 3.25)
    add_shape(slide, x, Inches(1.3), Inches(3.0), Inches(5.8),
              fill_color=CARD_BG, line_color=clr)
    add_text(slide, x + Inches(0.2), Inches(1.45), Inches(2.6), Inches(0.5),
             title, 20, clr, True, PP_ALIGN.CENTER, FONT)
    ornamental_rule(slide, x + Inches(0.2), Inches(1.95), Inches(2.6))
    add_bullet_list(slide, x + Inches(0.2), Inches(2.3), Inches(2.6), Inches(4.5),
                    items, 13, DARK_BROWN, "✦")


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — LIVE DEMO
# ═══════════════════════════════════════════════════════════════════════
section_divider(
    "LIVE DEMONSTRATION",
    "Manager's Command Center  ·  Analyze  ·  Chat  ·  Transcribe",
    "Rockwell-style: An excited crowd gathered around a store window display, "
    "watching a new product demo — representing the live demonstration"
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — MEET THE TEAM DIVIDER
# ═══════════════════════════════════════════════════════════════════════
section_divider(
    "MEET THE TEAM",
    "Individual Contributions",
    "Rockwell-style: Four craftspeople at a long workbench, each focused on their "
    "own specialty but working side by side — representing team collaboration"
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — TROY
# ═══════════════════════════════════════════════════════════════════════
dev_slide(
    name="Troy",
    role="CI/CD & Project Management",
    emoji="🔧",
    contributions=[
        "Designed and implemented all 6 GitHub Actions workflows",
        "terraform.yml — plan/apply with environment selection (dev/staging/prod)",
        "sagemaker-deploy.yml — endpoint lifecycle with health polling + inference tests",
        "deploy.yml — matrix Docker build → GHCR → EKS rollout + smoke tests",
        "ci-post-merge.yml — automated 4-job pipeline on merge to main",
        "slack-pr-events.yml — team notifications on PR open/merge",
        "teardown.yml — full infrastructure teardown workflow",
        "Set up GitHub Projects Kanban board for sprint tracking",
        "Branch protection rules and PR review workflow",
    ],
    tech_stack=[
        "GitHub Actions (YAML, matrix builds)",
        "Docker & GHCR (container registry)",
        "kubectl (EKS rollout management)",
        "Slack Webhooks",
        "GitHub Projects (Kanban)",
    ],
    talking_points=[
        "Walk through the CI/CD pipeline diagram",
        "Show the manual dispatch → auto validation flow",
        "Demonstrate Slack notification on PR merge",
        "Explain how post-merge pipeline catches regressions",
    ],
    img_desc="Rockwell-style: A determined mail carrier delivering packages through "
             "rain and snow — representing reliable CI/CD delivery",
    accent=NAVY,
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — KATHLEEN
# ═══════════════════════════════════════════════════════════════════════
dev_slide(
    name="Kathleen",
    role="ML Engineering — Churn Prediction & Frontend",
    emoji="📊",
    contributions=[
        "Trained XGBoost churn model v3: 95% accuracy, 0.9861 AUC",
        "Engineered 31 features including 7 cross-model Agent 1 features",
        "Deployed model to SageMaker endpoint (XGBoost container, native format)",
        "Built churn-predictor-api with internal S3 customer data lookup",
        "/predict, /customers, /customer-details, /high-risk endpoints",
        "Co-built LangGraph retention_agent.py with 4-tool orchestration",
        "Designed TriLink product catalog with risk-based retention actions",
        "Output guardrails for agent recommendations",
        "Built React + Tailwind frontend (Analyze, Chat, Transcribe tabs)",
    ],
    tech_stack=[
        "SageMaker (XGBoost, endpoint deployment)",
        "FastAPI + boto3 + pandas",
        "LangGraph + LangChain agents",
        "React 18 + TypeScript + Vite + Tailwind",
        "Docker (multi-stage builds)",
    ],
    talking_points=[
        "Walk through churn model training + feature engineering",
        "Show the 31-feature pipeline + Agent 1 enrichment",
        "Demo the /high-risk batch prediction endpoint",
        "Explain the LangGraph agent routing + guardrails",
    ],
    img_desc="Rockwell-style: A scientist peering through a microscope with intense "
             "focus, colorful charts on the wall behind — representing data analysis",
    accent=DEEP_RED,
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — OKINO
# ═══════════════════════════════════════════════════════════════════════
dev_slide(
    name="Okino",
    role="ML Engineering — Sentiment Analysis",
    emoji="🧠",
    contributions=[
        "Built sentiment-analysis-api with Bedrock-powered analysis",
        "Training notebook for QA/transcript evaluation model",
        "FastAPI wrapper returning qa_score, sentiment, emotion metrics",
        "Designed the 7 Agent 1 features that enrich churn predictions:",
        "  → sentiment, emotion_frustration, emotion_anger",
        "  → sentiment_shift, escalation_flag, resolution_flag, qa_score",
        "Co-built LangGraph agent tool contracts for cross-service calls",
        "Dockerfile + K8s deployment and service manifests",
        "SageMaker endpoint configuration in Terraform",
    ],
    tech_stack=[
        "SageMaker (model training + endpoints)",
        "Amazon Bedrock (Claude for sentiment)",
        "FastAPI + boto3",
        "Docker + Kubernetes manifests",
        "Terraform (sagemaker.tf)",
    ],
    talking_points=[
        "Explain how call transcripts become 7 structured features",
        "Show the sentiment → churn enrichment pipeline",
        "Walk through the /predict response schema",
        "Describe how Agent 1 features improve churn accuracy",
    ],
    img_desc="Rockwell-style: A wise judge listening carefully to testimony in a courtroom, "
             "weighing emotions and truth — representing sentiment analysis",
    accent=FOREST,
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — GEORGE
# ═══════════════════════════════════════════════════════════════════════
dev_slide(
    name="George",
    role="Infrastructure & Platform Engineering",
    emoji="⚙️",
    contributions=[
        "Authored all Terraform IaC: S3, IAM, SageMaker, EKS, Bedrock, Transcribe",
        "Designed & deployed full Kubernetes architecture (EKS cluster)",
        "All K8s manifests: deployments, services, configmaps, secrets, quotas",
        "Built agent-service: FastAPI + LangGraph /chat with session memory",
        "Integrated Bedrock Guardrails (PII anonymization, content filtering)",
        "docker-compose.yml for full local development stack",
        "Amazon Transcribe pipeline: S3 audio → Lambda → Transcribe → S3",
        "Terraform for Lambda + IAM + S3 event trigger",
    ],
    tech_stack=[
        "Terraform (14 config files, multi-resource)",
        "Amazon EKS + kubectl",
        "Kubernetes (namespace, RBAC, resource quotas)",
        "FastAPI + LangGraph + Bedrock",
        "AWS Lambda + Amazon Transcribe",
        "Docker + docker-compose",
    ],
    talking_points=[
        "Walk through the Terraform resource graph",
        "Show the K8s namespace layout and service mesh",
        "Explain the dual-LLM guardrail architecture",
        "Demo the Transcribe pipeline: audio → transcript",
    ],
    img_desc="Rockwell-style: A master plumber surrounded by pipes, valves, and blueprints, "
             "confidently connecting a complex system — representing infrastructure engineering",
    accent=WARM_GOLD,
)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — CI/CD PIPELINE
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, NAVY)
add_text(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7),
         "CI/CD PIPELINE — 6 GitHub Actions Workflows", 32, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

workflows = [
    ("terraform.yml", "Manual Dispatch", "Terraform plan → apply  (dev/staging/prod)", DEEP_RED),
    ("sagemaker-deploy.yml", "Manual Dispatch", "Deploy/delete endpoints + health poll + inference test", FOREST),
    ("deploy.yml", "Manual Dispatch", "Matrix Docker build → GHCR → EKS rollout + smoke test", NAVY),
    ("ci-post-merge.yml", "Push to main", "Unit tests → Docker → Endpoint health → E2E smoke", WARM_GOLD),
    ("slack-pr-events.yml", "PR open/merge", "Formatted Slack notifications to team channel", BURNT_ORANGE),
    ("teardown.yml", "Manual Dispatch", "Full infrastructure teardown", BROWN),
]
for i, (name, trigger, desc, clr) in enumerate(workflows):
    y = Inches(1.3 + i * 0.98)
    add_shape(slide, Inches(0.6), y, Inches(12.1), Inches(0.85),
              fill_color=CARD_BG, line_color=clr)
    add_text(slide, Inches(1.0), y + Inches(0.08), Inches(3.5), Inches(0.35),
             name, 17, clr, True, font_name=FONT)
    add_text(slide, Inches(1.0), y + Inches(0.45), Inches(3.5), Inches(0.35),
             trigger, 12, SLATE, font_name=FONT, italic=True)
    add_text(slide, Inches(4.8), y + Inches(0.2), Inches(7.5), Inches(0.5),
             desc, 16, DARK_BROWN, font_name=FONT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — AGENTIC AI
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, DEEP_RED)
add_text(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7),
         "AGENTIC AI — LangGraph Orchestration", 34, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

# Agent nodes
nodes = [
    ("DataGatherer", "Bedrock + Guardrails — Collects customer data, invokes tools", NAVY, Inches(1.4)),
    ("RetentionStrategist", "Clean LLM (no guardrail) — Generates retention offers from risk data", DEEP_RED, Inches(3.0)),
]
for name, desc, clr, y in nodes:
    add_shape(slide, Inches(0.6), y, Inches(5.8), Inches(1.35),
              fill_color=CARD_BG, line_color=clr)
    add_text(slide, Inches(1.0), y + Inches(0.1), Inches(5), Inches(0.4),
             name, 22, clr, True, font_name=FONT)
    add_text(slide, Inches(1.0), y + Inches(0.55), Inches(5.2), Inches(0.7),
             desc, 14, DARK_BROWN, font_name=FONT)

# Tools list
add_shape(slide, Inches(6.8), Inches(1.4), Inches(5.9), Inches(2.95),
          fill_color=CARD_BG, line_color=FOREST)
add_text(slide, Inches(7.1), Inches(1.55), Inches(5.2), Inches(0.4),
         "5 Agent Tools", 22, FOREST, True, font_name=FONT)
tools = [
    "get_customer_details → Churn API /customer-details/{id}",
    "predict_churn → Churn API /predict (31 features → SageMaker)",
    "analyze_call → Sentiment API /predict (transcript → 7 metrics)",
    "get_high_risk_customers → Churn API /high-risk (batch)",
    "get_transcripts → S3 transcript listing per customer",
]
add_bullet_list(slide, Inches(7.1), Inches(2.1), Inches(5.5), Inches(2.3), tools, 13, DARK_BROWN, "✦")

# Guardrails
add_shape(slide, Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.5),
          fill_color=CARD_BG, line_color=WARM_GOLD)
add_text(slide, Inches(1.0), Inches(4.75), Inches(5), Inches(0.4),
         "Bedrock Guardrails", 20, WARM_GOLD, True, font_name=FONT)
guardrails = [
    "Hate/Sexual content: HIGH filter  ·  Violence: HIGH/MEDIUM  ·  Misconduct: HIGH/LOW",
    "PII Protection: ANONYMIZE names, emails, phone, SSN  ·  BLOCK credit card, bank account",
    "Dual-LLM design: DataGatherer uses guardrails; RetentionStrategist uses clean LLM to avoid blocking domain terms",
]
add_bullet_list(slide, Inches(1.0), Inches(5.3), Inches(11.5), Inches(1.5), guardrails, 14, DARK_BROWN, "✦")

# Side illustration
add_image_placeholder(slide, Inches(0.6), Inches(4.65), Inches(0), Inches(0), "", WARM_GOLD)  # skip, too crowded


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — KEY METRICS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, FOREST)
add_text(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7),
         "KEY METRICS & RESULTS", 36, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

metrics = [
    ("95%", "Model\nAccuracy", FOREST),
    ("0.986", "AUC\nScore", NAVY),
    ("31", "Engineered\nFeatures", DEEP_RED),
    ("5", "Agent\nTools", WARM_GOLD),
    ("3", "Micro-\nservices", BURNT_ORANGE),
    ("6", "CI/CD\nWorkflows", BROWN),
]
for i, (val, label, clr) in enumerate(metrics):
    x = Inches(0.4 + i * 2.15)
    add_shape(slide, x, Inches(1.3), Inches(1.9), Inches(2.2),
              fill_color=CARD_BG, line_color=clr)
    add_text(slide, x + Inches(0.1), Inches(1.5), Inches(1.7), Inches(1),
             val, 44, clr, True, PP_ALIGN.CENTER, FONT)
    add_text(slide, x + Inches(0.1), Inches(2.5), Inches(1.7), Inches(0.8),
             label, 14, SLATE, False, PP_ALIGN.CENTER, FONT)

highlights = [
    "End-to-end ML pipeline: raw data → training → SageMaker → FastAPI wrapper → agent tool → user",
    "Agentic AI with multi-step reasoning: transcript → sentiment → churn enrichment → retention strategy",
    "Full IaC: 14 Terraform files provisioning S3, IAM, SageMaker, EKS, Bedrock, Transcribe, Lambda",
    "Production K8s: namespace isolation, resource quotas, health probes, ConfigMaps, rolling deploys",
    "Amazon Transcribe: audio upload → speaker diarization → structured transcript → S3 storage",
]
add_shape(slide, Inches(0.4), Inches(3.8), Inches(12.5), Inches(3.3),
          fill_color=CARD_BG, line_color=NAVY)
add_text(slide, Inches(0.7), Inches(3.95), Inches(5), Inches(0.4),
         "Platform Highlights", 22, NAVY, True, font_name=FONT)
add_bullet_list(slide, Inches(0.7), Inches(4.5), Inches(12), Inches(2.8), highlights, 15, DARK_BROWN, "✦")


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — LESSONS LEARNED
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, BROWN)
add_text(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7),
         "LESSONS LEARNED", 36, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

lessons_left = [
    "Define API contracts early — enabled parallel development",
    "Bedrock Guardrails block domain terms (cancel, churn) — needed dual-LLM workaround",
    "SageMaker cold starts affect initial latency — added health polling in CI",
    "Multi-stage Docker builds cut image sizes significantly",
]
lessons_right = [
    "LangGraph > raw LangChain for multi-step agent workflows",
    "docker-compose essential for local integration testing before K8s",
    "Slack notifications keep the entire team informed on PR flow",
    "Feature engineering (31 features) matters more than model choice",
]

add_shape(slide, Inches(0.6), Inches(1.3), Inches(5.8), Inches(4.0),
          fill_color=CARD_BG, line_color=DEEP_RED)
add_text(slide, Inches(0.9), Inches(1.45), Inches(5.2), Inches(0.4),
         "Technical Insights", 22, DEEP_RED, True, font_name=FONT)
ornamental_rule(slide, Inches(0.9), Inches(1.9), Inches(5.2))
add_bullet_list(slide, Inches(0.9), Inches(2.25), Inches(5.2), Inches(3.0), lessons_left, 15, DARK_BROWN, "✦")

add_shape(slide, Inches(6.8), Inches(1.3), Inches(5.9), Inches(4.0),
          fill_color=CARD_BG, line_color=FOREST)
add_text(slide, Inches(7.1), Inches(1.45), Inches(5.2), Inches(0.4),
         "Process Insights", 22, FOREST, True, font_name=FONT)
ornamental_rule(slide, Inches(7.1), Inches(1.9), Inches(5.2))
add_bullet_list(slide, Inches(7.1), Inches(2.25), Inches(5.2), Inches(3.0), lessons_right, 15, DARK_BROWN, "✦")

# Bottom illustration
add_image_placeholder(slide, Inches(3.5), Inches(5.5), Inches(6.3), Inches(1.7),
    "Rockwell-style: A student at a desk surrounded by open books, having an 'aha!' moment — "
    "representing lessons learned",
    BROWN)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 16 — FUTURE IMPROVEMENTS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, WARM_GOLD)
add_text(slide, Inches(0.8), Inches(0.15), Inches(11.5), Inches(0.7),
         "FUTURE IMPROVEMENTS", 36, CREAM_TEXT, True, PP_ALIGN.LEFT, FONT)

# ── Card 1: Historical Sentiment Trending ──
add_shape(slide, Inches(0.5), Inches(1.2), Inches(12.3), Inches(1.95),
          fill_color=CARD_BG, line_color=DEEP_RED)
add_text(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.45),
         "Historical Sentiment Trending & Cross-Call Analysis", 22, DEEP_RED, True, font_name=FONT)
add_text(slide, Inches(0.8), Inches(1.75), Inches(5.5), Inches(0.35),
         "The Problem", 14, WARM_GOLD, True, font_name=FONT)
add_text(slide, Inches(0.8), Inches(2.05), Inches(5.5), Inches(0.8),
         "Each call is analyzed in isolation — a manager sees one snapshot. But the real signal is "
         "the trajectory. A customer trending from positive to negative over three months is a very "
         "different risk than a first-time angry caller with a billing error.",
         12, DARK_BROWN, font_name=FONT)
add_text(slide, Inches(6.8), Inches(1.75), Inches(5.5), Inches(0.35),
         "Future Implementation", 14, FOREST, True, font_name=FONT)
add_text(slide, Inches(6.8), Inches(2.05), Inches(5.7), Inches(0.8),
         "Auto-pull all historical transcripts from S3, run sentiment enrichment on each, and present "
         "a sentiment timeline. Incorporate features like \"calls in last 90 days,\" \"average sentiment "
         "trend,\" and \"worst sentiment delta\" — serving a longitudinal risk model.",
         12, DARK_BROWN, font_name=FONT)

# ── Card 2: Persistent Memory ──
add_shape(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(1.95),
          fill_color=CARD_BG, line_color=NAVY)
add_text(slide, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.45),
         "Persistent Conversation Memory with Semantic Search", 22, NAVY, True, font_name=FONT)
add_text(slide, Inches(0.8), Inches(3.85), Inches(5.5), Inches(0.35),
         "The Problem", 14, WARM_GOLD, True, font_name=FONT)
add_text(slide, Inches(0.8), Inches(4.15), Inches(5.5), Inches(0.8),
         "Agent memory resets on service restart. A manager who discussed a customer yesterday "
         "must re-explain context today. No way to search past conversations — \"what did we "
         "decide about C00036458 last week?\" has no answer.",
         12, DARK_BROWN, font_name=FONT)
add_text(slide, Inches(6.8), Inches(3.85), Inches(5.5), Inches(0.35),
         "Future Implementation", 14, FOREST, True, font_name=FONT)
add_text(slide, Inches(6.8), Inches(4.15), Inches(5.7), Inches(0.8),
         "Replace MemorySaver with PostgreSQL + pgvector (or Amazon MemoryDB). Store conversation "
         "history with embeddings for semantic search — find related discussions across all sessions. "
         "\"Have we offered loyalty discounts to any Basic_25 customers this quarter?\"",
         12, DARK_BROWN, font_name=FONT)

# ── Card 3: Dashboard & Outcome Tracking ──
add_shape(slide, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.95),
          fill_color=CARD_BG, line_color=FOREST)
add_text(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.45),
         "Retention Action Dashboard & Outcome Tracking", 22, FOREST, True, font_name=FONT)
add_text(slide, Inches(0.8), Inches(5.95), Inches(5.5), Inches(0.35),
         "The Problem", 14, WARM_GOLD, True, font_name=FONT)
add_text(slide, Inches(0.8), Inches(6.25), Inches(5.5), Inches(0.8),
         "No aggregate view — how many high-risk customers exist today? What actions were "
         "recommended this week? Did they work? Without outcome tracking, the retention engine "
         "is a recommendation tool with no feedback loop.",
         12, DARK_BROWN, font_name=FONT)
add_text(slide, Inches(6.8), Inches(5.95), Inches(5.5), Inches(0.35),
         "Future Implementation", 14, FOREST, True, font_name=FONT)
add_text(slide, Inches(6.8), Inches(6.25), Inches(5.7), Inches(0.8),
         "Add a Dashboard tab with aggregate metrics: risk distribution, actions taken, trend charts. "
         "Log retention actions, compare sentiment before/after. Feed outcomes back as reinforcement — "
         "the system learns which actions actually reduce churn for which customer profiles.",
         12, DARK_BROWN, font_name=FONT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 17 — AI USAGE TRENDS
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
top_banner(slide, NAVY)
add_text(slide, Inches(0.8), Inches(0.1), Inches(11.5), Inches(0.7),
         "HOW WE USED AI  —  4 Key Trends", 38, CREAM_TEXT, True, PP_ALIGN.CENTER, FONT)

ai_trends = [
    ("Debugging Complex Systems Was AI's Strongest Role",
     "AI was most impactful diagnosing failures across interconnected services — tracing "
     "CI/CD workflow failures through CloudWatch logs (Troy), eliminating SageMaker deploy "
     "root causes across 5 failed attempts (Kathleen), and separating timeout vs HTTP error "
     "handling in LangChain tools (George). AI read logs and proposed targeted fixes faster "
     "than manual documentation searches.",
     DEEP_RED),
    ("Human Direction Turned AI Execution Into Quality",
     "The most impactful moments were when team members redirected AI rather than accepting "
     "its first answer. Kathleen's pivot — \"evaluate every possible issue with probability "
     "of success\" — broke a debugging loop and led to the working XGBoost solution. Troy "
     "told the AI to \"look it up from deploy.py, don't guess\" to prevent duplicate "
     "$400/mo endpoints. AI executed capably but needed human judgment and domain knowledge.",
     FOREST),
    ("Infrastructure & DevOps Dominated AI Workload",
     "The heaviest AI usage was Terraform definitions, GitHub Actions workflows, SageMaker "
     "deployment scripts, K8s health probes, and Docker builds — configuration-heavy work "
     "with strict syntax requirements. Troy's 58-task log is essentially CI/CD pipeline "
     "implementation. AI handled boilerplate-rich infrastructure well because it could "
     "synthesize scattered documentation into working configurations.",
     WARM_GOLD),
    ("AI Enabled Rapid Iteration Across 10+ AWS Services",
     "The project spanned SageMaker, EKS, Bedrock, Lambda, Transcribe, S3, IAM, CloudWatch, "
     "ECR, and DynamoDB — each with service-specific edge cases. AI synthesized knowledge "
     "across these boundaries: kubeconfig ARN quirks breaking CI, GitHub Actions setting "
     "missing secrets to empty strings, SageMaker rollbacks silently serving stale models. "
     "This cross-service synthesis was a force multiplier.",
     BURNT_ORANGE),
]

card_top = Inches(1.3)
card_height = Inches(1.4)
card_gap = Inches(0.15)
for i, (title, body, accent) in enumerate(ai_trends):
    y = card_top + i * (card_height + card_gap)
    # accent stripe on left edge
    add_rect(slide, Inches(0.5), y, Inches(0.08), card_height,
             fill_color=accent, line_color=None)
    # card background
    add_shape(slide, Inches(0.58), y, Inches(12.2), card_height,
              fill_color=CARD_BG, line_color=BROWN, line_width=Pt(1))
    # trend number
    add_text(slide, Inches(0.75), y + Inches(0.1), Inches(0.5), Inches(0.4),
             str(i + 1), 28, accent, True, PP_ALIGN.CENTER, FONT)
    # title
    add_text(slide, Inches(1.3), y + Inches(0.08), Inches(11), Inches(0.35),
             title, 18, accent, True, PP_ALIGN.LEFT, FONT)
    # body
    add_text(slide, Inches(1.3), y + Inches(0.45), Inches(11.2), Inches(0.9),
             body, 11, DARK_BROWN, False, PP_ALIGN.LEFT, FONT)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 18 — Q&A / THANK YOU
# ═══════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, PARCHMENT)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(2.6),
         fill_color=NAVY, line_color=None)
add_rect(slide, Inches(0), Inches(2.56), Inches(13.333), Inches(0.06),
         fill_color=WARM_GOLD, line_color=None)

add_text(slide, Inches(1), Inches(0.5), Inches(11.3), Inches(1.3),
         "QUESTIONS?", 62, CREAM_TEXT, True, PP_ALIGN.CENTER, FONT)
add_text(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(0.6),
         "THE RETENTION ENGINE", 28, WARM_GOLD, False, PP_ALIGN.CENTER, FONT)

add_image_placeholder(slide, Inches(4.2), Inches(3.0), Inches(4.9), Inches(2.8),
    "Rockwell-style: Four friends gathered around a table, warmly shaking hands after "
    "completing a big project — representing team accomplishment",
    NAVY)

ornamental_rule(slide, Inches(2.5), Inches(6.0), Inches(8.3))

add_text(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
         "Troy  ·  Kathleen  ·  Okino  ·  George", 24, NAVY, True, PP_ALIGN.CENTER, FONT)
add_text(slide, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
         "github.com/Lumin33r/Capstone-Churn", 16, SEPIA, False, PP_ALIGN.CENTER, FONT)


# ── Save ────────────────────────────────────────────────────────────────
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "TheRetentionEngine_Rockwell.pptx")
prs.save(out_path)
print(f"✅ Saved → {out_path}")
